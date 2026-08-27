import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Light Theme Colors ──
    WHITE = RGBColor(255, 255, 255)
    OFF_WHITE = RGBColor(248, 250, 252)
    TEAL = RGBColor(13, 148, 136)
    INDIGO = RGBColor(49, 46, 129)
    DARK = RGBColor(31, 41, 55)
    GRAY = RGBColor(107, 114, 128)
    LIGHT_GRAY = RGBColor(229, 231, 235)
    CORAL = RGBColor(239, 68, 68)
    EMERALD = RGBColor(16, 185, 129)
    PURPLE = RGBColor(124, 58, 237)
    BLUE = RGBColor(59, 130, 246)
    PINK = RGBColor(236, 72, 153)
    AMBER = RGBColor(245, 158, 11)

    # Image paths
    LOGO = r"C:\Users\Yash\.gemini\antigravity\brain\0c9e94ab-367d-4800-bea6-9239bfe78090\mockmirror_logo_light_1786113030507.jpg"
    WORKFLOW = r"C:\Users\Yash\.gemini\antigravity\brain\0c9e94ab-367d-4800-bea6-9239bfe78090\.user_uploaded\media_1786112946054.jpg"
    FRONTEND_LOGOS = r"C:\Users\Yash\.gemini\antigravity\brain\0c9e94ab-367d-4800-bea6-9239bfe78090\frontend_logos_v2_1786118526971.jpg"
    BACKEND_LOGOS = r"C:\Users\Yash\.gemini\antigravity\brain\0c9e94ab-367d-4800-bea6-9239bfe78090\backend_tech_logos_1786117088954.jpg"

    def set_bg(slide, color=WHITE):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def accent_bar(slide, color=TEAL):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

    def bottom_bar(slide, color=INDIGO):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.32), prs.slide_width, Inches(0.18))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

    def divider(slide, x, y, w, color=TEAL):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.035))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

    def txt(slide, l, t, w, h, text, fname="Segoe UI", sz=16, bold=False, italic=False, color=DARK, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text
        p.font.name = fname; p.font.size = Pt(sz); p.font.bold = bold
        p.font.italic = italic; p.font.color.rgb = color; p.alignment = align
        return tf

    def bullet_block(slide, l, t, w, h, items, ts=16, ds=15, tc=TEAL, dc=DARK):
        """items = list of (bold_title, description)"""
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for title, desc in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(10); p.space_after = Pt(3)
            r1 = p.add_run()
            r1.text = f"•  {title}:  "
            r1.font.name = "Segoe UI"; r1.font.size = Pt(ts)
            r1.font.bold = True; r1.font.color.rgb = tc
            r2 = p.add_run()
            r2.text = desc
            r2.font.name = "Segoe UI"; r2.font.size = Pt(ds)
            r2.font.bold = False; r2.font.color.rgb = dc
        return tf

    def simple_bullets(slide, l, t, w, h, items, sz=14, color=DARK, bc=TEAL, indent_items=None):
        """indent_items = set of indices that should be indented sub-bullets"""
        if indent_items is None:
            indent_items = set()
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(4); p.space_after = Pt(2)
            prefix = "     •  " if i in indent_items else "•  "
            rb = p.add_run(); rb.text = prefix
            rb.font.name = "Segoe UI"; rb.font.size = Pt(sz if i not in indent_items else sz - 1)
            rb.font.color.rgb = bc; rb.font.bold = True
            rt = p.add_run(); rt.text = item
            rt.font.name = "Segoe UI"; rt.font.size = Pt(sz if i not in indent_items else sz - 1)
            rt.font.color.rgb = color
        return tf

    # ═══════════════════════════════════════════════════════
    # SLIDE 1: TITLE — MockMirror : Practice. Reflect. Improve.
    # ═══════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s1)
    accent_bar(s1, TEAL)
    bottom_bar(s1, INDIGO)

    # Left teal-tinted panel
    lp = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.06), Inches(5.8), Inches(7.26))
    lp.fill.solid(); lp.fill.fore_color.rgb = RGBColor(240, 253, 250); lp.line.fill.background()

    try:
        s1.shapes.add_picture(LOGO, Inches(0.7), Inches(0.6), Inches(4.4), Inches(2.5))
    except: pass

    txt(s1, 0.7, 3.3, 4.8, 0.7, "MockMirror", "Segoe UI", 44, True, False, INDIGO)
    txt(s1, 0.7, 3.95, 4.8, 0.5, "Practice. Reflect. Improve.", "Segoe UI", 20, False, True, TEAL)
    txt(s1, 0.7, 4.5, 4.8, 0.4, "AI-Powered Multimodal Interview Simulator", "Segoe UI", 13, False, False, GRAY)

    # Right — Team Name + Members
    txt(s1, 6.5, 1.0, 6.0, 0.4, "TEAM  TERMINAL4", "Segoe UI", 18, True, False, INDIGO)
    divider(s1, 6.5, 1.4, 5.0, TEAL)
    txt(s1, 6.5, 1.5, 6.0, 0.35, "TEAM MEMBERS", "Segoe UI", 13, True, False, GRAY)

    members = [
        ("Yash Prajapati", "Enrollment: 250413116015"),
        ("Tisha Kotak", "Enrollment: 240410116132"),
        ("Yesha Patel", "Enrollment: 250413116015"),
        ("Hetvi Shah", "Enrollment: 250413107019"),
    ]
    for i, (name, enroll) in enumerate(members):
        y = 2.0 + (i * 1.15)
        badge = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(y), Inches(0.42), Inches(0.42))
        badge.fill.solid(); badge.fill.fore_color.rgb = TEAL; badge.line.fill.background()
        badge.text_frame.paragraphs[0].text = str(i + 1)
        badge.text_frame.paragraphs[0].font.name = "Segoe UI"
        badge.text_frame.paragraphs[0].font.size = Pt(14)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf = txt(s1, 7.15, y - 0.02, 5.0, 0.5, name, "Segoe UI", 16, True, False, DARK)
        p2 = tf.add_paragraph()
        p2.text = enroll; p2.font.name = "Segoe UI"; p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.LEFT

    # ═══════════════════════════════════════════════════════
    # SLIDE 2: THE PROBLEM (exact content from screenshot)
    # ═══════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s2)
    accent_bar(s2, CORAL)
    bottom_bar(s2, INDIGO)

    txt(s2, 1.0, 0.35, 11.33, 0.7, "THE PROBLEM", "Segoe UI", 38, True, False, INDIGO)
    txt(s2, 1.0, 0.95, 11.33, 0.35, "Why Interview Preparation is Broken", "Segoe UI", 15, False, True, GRAY)
    divider(s2, 1.0, 1.35, 2.0, CORAL)

    bullet_block(s2, 1.0, 1.6, 11.33, 5.5, [
        ("Placement & Job Interview Anxiety",
         "Students and freshers face extreme nervousness during campus placements and job interviews due to lack of real practice with a live interviewer."),
        ("Hesitation While Speaking",
         "Many candidates know the answers but freeze, stutter, or hesitate when speaking out loud, leading to poor first impressions."),
        ("Unknown Question Patterns",
         "Candidates have no idea what type of questions will be asked — technical, behavioral, situational, or resume-based — and cannot prepare effectively."),
        ("No Feedback on Body Language",
         "Traditional preparation ignores facial expressions, eye contact, posture, and head movement — all critical factors interviewers silently evaluate."),
        ("Static Mock Interviews",
         "Existing tools provide scripted, pre-written Q&As that don't adapt to a candidate's actual resume, spoken answers, or real-time performance."),
        ("No Safe Practice Environment",
         "There is no affordable, private, judgment-free space where candidates can repeatedly practice and improve before the real interview day."),
    ], 16, 15, CORAL, DARK)

    # ═══════════════════════════════════════════════════════
    # SLIDE 3: OUR SOLUTION — MOCKMIRROR (exact content)
    # ═══════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s3)
    accent_bar(s3, EMERALD)
    bottom_bar(s3, INDIGO)

    txt(s3, 1.0, 0.35, 11.33, 0.7, "OUR SOLUTION — MOCKMIRROR", "Segoe UI", 38, True, False, INDIGO)
    txt(s3, 1.0, 0.95, 11.33, 0.35, "Practice Smart. Interview Confident.", "Segoe UI", 15, False, True, TEAL)
    divider(s3, 1.0, 1.35, 2.0, EMERALD)

    bullet_block(s3, 1.0, 1.6, 11.33, 5.5, [
        ("AI-Powered Conversational Practice",
         "Candidates visit our website and practice for their interviews by having a real-time voice conversation with our AI interviewer — just like a real interview."),
        ("Resume-Based Personalized Questions",
         "Upload your PDF resume, and the AI reads your skills, projects, and experience to ask highly targeted and relevant questions."),
        ("Real-Time Answer Evaluation",
         "The AI doesn't just ask questions — it analyzes every answer for correctness, depth, and relevance. If an answer is fake or vague, it immediately calls it out and corrects the candidate."),
        ("Face Expression & Posture Detection",
         "Using the browser camera, MockMirror tracks facial expressions, face centering, eye contact, and head posture — providing feedback on how the candidate physically presents themselves."),
        ("Comprehensive Performance Report",
         "A detailed report is generated with scores across technical accuracy, communication, confidence, and behavioral metrics — downloadable as a PDF."),
    ], 16, 15, EMERALD, DARK)

    # ═══════════════════════════════════════════════════════
    # SLIDE 4: USER WORKFLOW (workflow diagram image)
    # ═══════════════════════════════════════════════════════
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s4)
    accent_bar(s4, BLUE)
    bottom_bar(s4, INDIGO)

    txt(s4, 1.0, 0.3, 11.33, 0.55, "USER WORKFLOW", "Segoe UI", 38, True, False, INDIGO)
    txt(s4, 1.0, 0.8, 11.33, 0.35, "How a Candidate Uses MockMirror — Step by Step", "Segoe UI", 15, False, True, GRAY)
    divider(s4, 1.0, 1.18, 2.0, BLUE)

    try:
        s4.shapes.add_picture(WORKFLOW, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.8))
    except Exception as e:
        print(f"Workflow image error: {e}")

    # ═══════════════════════════════════════════════════════
    # SLIDE 5: FUTURE SCOPE (exact content from screenshot)
    # ═══════════════════════════════════════════════════════
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s5)
    accent_bar(s5, PURPLE)
    bottom_bar(s5, INDIGO)

    txt(s5, 1.0, 0.35, 11.33, 0.7, "FUTURE SCOPE", "Segoe UI", 38, True, False, INDIGO)
    txt(s5, 1.0, 0.95, 11.33, 0.35, "Scaling MockMirror Beyond Candidate Self-Practice", "Segoe UI", 15, False, True, GRAY)
    divider(s5, 1.0, 1.35, 2.0, PURPLE)

    # Current Version Card
    c1 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(2.8))
    c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(240, 253, 250)
    c1.line.color.rgb = EMERALD; c1.line.width = Pt(1.5)

    txt(s5, 1.1, 1.75, 5.0, 0.35, "CURRENT VERSION (Candidate Mode)", "Segoe UI", 14, True, False, EMERALD)

    bullet_block(s5, 1.1, 2.15, 5.0, 2.0, [
        ("Self-Practice",
         "Candidate uploads resume, selects domain & difficulty, and practices independently with AI."),
        ("AI-Generated Questions",
         "Questions are auto-generated based on resume, domain, and difficulty level."),
        ("Behavioral Tracking",
         "Face detection and posture tracking provide visual behavior feedback."),
    ], 14, 13, EMERALD, DARK)

    # Future Version Card
    c2 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.8))
    c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(253, 242, 248)
    c2.line.color.rgb = PINK; c2.line.width = Pt(1.5)

    txt(s5, 7.1, 1.75, 5.2, 0.35, "FUTURE VERSION (Interviewer Mode)", "Segoe UI", 14, True, False, PINK)

    bullet_block(s5, 7.1, 2.15, 5.2, 2.0, [
        ("Interviewer Dashboard",
         "The interviewer selects topics and uploads custom data/syllabus for the AI to prepare questions from."),
        ("AI Conducts on Behalf",
         "The AI interviews the candidate on behalf of the interviewer using the provided material."),
        ("Automated Candidate Reports",
         "The interviewer receives a detailed performance report of each candidate without being present."),
    ], 14, 13, PINK, DARK)

    # Additional Future Enhancements
    txt(s5, 1.0, 4.7, 11.33, 0.35, "ADDITIONAL FUTURE ENHANCEMENTS", "Segoe UI", 14, True, False, GRAY)
    divider(s5, 1.0, 5.05, 1.5, PURPLE)

    bullet_block(s5, 1.0, 5.15, 11.33, 2.0, [
        ("Sentiment & Stress Analysis",
         "Detect micro-expressions to assess confidence and stress in real-time."),
        ("Live Coding Whiteboard",
         "Integrate browser-based coding environments for technical screening."),
        ("Multi-Language & Accent Coaching",
         "Support regional languages and accent improvement modules."),
        ("LMS & Placement Cell Integration",
         "License MockMirror to universities and corporate HR portals."),
    ], 14, 13, PURPLE, DARK)

    # ═══════════════════════════════════════════════════════
    # SLIDE 6: TECHNOLOGY STACK (with logo images)
    # ═══════════════════════════════════════════════════════
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s6)
    accent_bar(s6, TEAL)
    bottom_bar(s6, INDIGO)

    txt(s6, 1.0, 0.35, 11.33, 0.7, "TECHNOLOGY STACK", "Segoe UI", 38, True, False, INDIGO)
    txt(s6, 1.0, 0.95, 11.33, 0.35, "Built with Modern, Privacy-First Technologies", "Segoe UI", 15, False, True, GRAY)
    divider(s6, 1.0, 1.35, 2.0, TEAL)

    # Frontend card
    fc = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(5.8), Inches(4.2))
    fc.fill.solid(); fc.fill.fore_color.rgb = RGBColor(240, 253, 250)
    fc.line.color.rgb = TEAL; fc.line.width = Pt(1.5)

    txt(s6, 1.0, 1.75, 5.0, 0.35, "FRONTEND", "Segoe UI", 18, True, False, TEAL)
    divider(s6, 1.0, 2.1, 1.2, TEAL)

    try:
        s6.shapes.add_picture(FRONTEND_LOGOS, Inches(0.8), Inches(2.25), Inches(5.4), Inches(3.3))
    except Exception as e:
        print(f"Frontend logos error: {e}")

    # Backend card
    bc = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.8), Inches(4.2))
    bc.fill.solid(); bc.fill.fore_color.rgb = RGBColor(238, 242, 255)
    bc.line.color.rgb = BLUE; bc.line.width = Pt(1.5)

    txt(s6, 7.3, 1.75, 5.0, 0.35, "BACKEND", "Segoe UI", 18, True, False, BLUE)
    divider(s6, 7.3, 2.1, 1.2, BLUE)

    try:
        s6.shapes.add_picture(BACKEND_LOGOS, Inches(7.1), Inches(2.25), Inches(5.4), Inches(3.3))
    except Exception as e:
        print(f"Backend logos error: {e}")

    # Architecture principles bar at bottom
    arch = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.05), Inches(12.1), Inches(1.1))
    arch.fill.solid(); arch.fill.fore_color.rgb = OFF_WHITE
    arch.line.color.rgb = LIGHT_GRAY; arch.line.width = Pt(1)

    txt(s6, 0.8, 6.1, 11.33, 0.3, "KEY ARCHITECTURE PRINCIPLES", "Segoe UI", 13, True, False, INDIGO)
    divider(s6, 0.8, 6.38, 1.5, TEAL)

    txt(s6, 0.8, 6.5, 3.5, 0.55, "🔒  100% Privacy\nNo database. No login. Resume parsed in-memory.", "Segoe UI", 11, False, False, DARK)
    txt(s6, 4.8, 6.5, 3.5, 0.55, "📷  Client-Side Vision\nAll camera processing runs locally in the browser.", "Segoe UI", 11, False, False, DARK)
    txt(s6, 8.8, 6.5, 3.5, 0.55, "⚡  Zero Friction\nOpen website → Upload resume → Start. Under 60s.", "Segoe UI", 11, False, False, DARK)

    # ═══════════════════════════════════════════════════════
    # SLIDE 7: THANK YOU / Q&A
    # ═══════════════════════════════════════════════════════
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s7)
    accent_bar(s7, TEAL)
    bottom_bar(s7, INDIGO)

    try:
        s7.shapes.add_picture(LOGO, Inches(4.2), Inches(0.5), Inches(5.0), Inches(2.8))
    except: pass

    txt(s7, 1.0, 3.5, 11.33, 0.8, "Thank You!", "Segoe UI", 52, True, False, INDIGO, PP_ALIGN.CENTER)
    txt(s7, 1.0, 4.3, 11.33, 0.5, "MockMirror", "Segoe UI", 28, True, False, TEAL, PP_ALIGN.CENTER)
    txt(s7, 1.0, 4.8, 11.33, 0.4, "Practice. Reflect. Improve.", "Segoe UI", 18, False, True, GRAY, PP_ALIGN.CENTER)
    divider(s7, 5.5, 5.35, 2.3, TEAL)
    txt(s7, 1.0, 5.5, 11.33, 0.5, "We are ready for your Questions & Answers", "Segoe UI", 18, False, False, DARK, PP_ALIGN.CENTER)

    # Save
    out = r"D:\M\r\MockMirror_Terminal4_v3.pptx"
    prs.save(out)
    print(f"Successfully created '{out}'!")

if __name__ == "__main__":
    create_presentation()
